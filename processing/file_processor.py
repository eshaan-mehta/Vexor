"""
File processing functions for the queue-based architecture.

This module contains the actual file processing logic that worker threads
will use to process tasks from the queue.
"""

import os
import hashlib
import mimetypes
import chromadb
import time
from datetime import datetime
from dataclasses import asdict
from typing import Optional
from chromadb.config import Settings
from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction
from pathlib import Path
import logging

# File content extraction imports
import chardet
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
import markdown

from models.filemetadata import FileMetadata
from .file_processing_queue import FileTask, TaskType
from enum import Enum

logger = logging.getLogger(__name__)

class ProcessingStatus(Enum):
    """Status of file processing operations."""
    SUCCESS = "success"
    SKIPPED = "skipped"
    HIDDEN = "hidden"
    LARGE = "large"
    FAILURE = "failure"

class FileProcessor:
    """
    Handles the actual processing of individual files.
    
    This class contains the logic for indexing, deleting, and moving files
    that worker threads will use to process tasks from the queue.
    """
    
    def __init__(self, db_path: str = "./chroma", 
                 metadata_collection_name: str = "file_metadata",
                 content_collection_name: str = "file_content"):
        """
        Initialize the file processor.
        
        Args:
            db_path: Path to ChromaDB database
            metadata_collection_name: Name of metadata collection
            content_collection_name: Name of content collection
        """
        self.db_path = db_path
        
        # Initialize ChromaDB client and collections
        os.makedirs(db_path, exist_ok=True)
        self.client = chromadb.PersistentClient(
            path=db_path,
            settings=Settings(anonymized_telemetry=False)
        )
        
        # Create embedding function per instance to avoid sharing resources
        self.embedding_function = SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")
        
        self.content_collection = self.client.get_or_create_collection(
            name=content_collection_name,
            embedding_function=self.embedding_function,
            metadata = {
                "hnsw:space": "cosine",  # cosine better for text embeddings
                "hnsw:construction_ef": 100,  # size of candidates during indexing (default = 100)
                "hnsw:search_ef": 100,  # size of candidates during searching (default = 100)
                "hnsw:M": 16  # max neighbours in node graph (default = 16)
            }
        )
        
        self.metadata_collection = self.client.get_or_create_collection(
            name=metadata_collection_name,
            embedding_function=self.embedding_function,
            metadata = {
                "hnsw:space": "cosine",  # cosine better for text embeddings
                "hnsw:construction_ef": 100,  # size of candidates during indexing (default = 100)
                "hnsw:search_ef": 100,  # size of candidates during searching (default = 100)
                "hnsw:M": 16  # max neighbours in node graph (default = 16)
            }
        )
        
        logger.info("FileProcessor initialized")
    
    def cleanup(self):
        """Clean up resources to prevent semaphore leaks."""
        try:
            # Close ChromaDB client if it has a close method
            if hasattr(self.client, 'close'):
                self.client.close()
            
            # Clean up embedding function if it has cleanup methods
            if hasattr(self.embedding_function, 'cleanup'):
                self.embedding_function.cleanup()
            elif hasattr(self.embedding_function, 'close'):
                self.embedding_function.close()
                
            logger.debug("FileProcessor resources cleaned up")
        except Exception as e:
            logger.error(f"Error during FileProcessor cleanup: {e}")
    
    def __del__(self):
        """Destructor to ensure cleanup on garbage collection."""
        self.cleanup()
    
    def process_task(self, task: FileTask) -> ProcessingStatus:
        """
        Process a file task based on its type.
        
        Args:
            task: FileTask to process
            
        Returns:
            ProcessingStatus: Status of the processing operation
        """
        try:
            if task.task_type == TaskType.INDEX_FILE or task.task_type == TaskType.UPDATE_FILE:
                return self._index_file(task.file_path)
            elif task.task_type == TaskType.DELETE_FILE:
                return ProcessingStatus.SUCCESS if self._delete_file(task.file_path) else ProcessingStatus.FAILURE
            elif task.task_type == TaskType.MOVE_FILE:
                return ProcessingStatus.SUCCESS if self._move_file(task) else ProcessingStatus.FAILURE
            else:
                logger.error(f"Unknown task type: {task.task_type}")
                return ProcessingStatus.FAILURE
                
        except Exception as e:
            logger.error(f"Error processing task {task.task_type} for {task.file_path}: {e}")
            return ProcessingStatus.FAILURE
    
    def _index_file(self, file_path: str) -> ProcessingStatus:
        """Index a single file."""
        if not os.path.exists(file_path) or os.path.isdir(file_path):
            logger.warning(f"File does not exist or is directory: {file_path}")
            return ProcessingStatus.FAILURE
        
        # Skip hidden files and temporary files
        name = os.path.basename(file_path)
        if name.startswith((".", "__", "~$")):
            logger.debug(f"Skipping hidden/temporary file: {name}")
            return ProcessingStatus.HIDDEN
        
        # Check file size limits
        if not self._check_file_size(file_path):
            return ProcessingStatus.LARGE
        
        try:
            # Extract metadata
            metadata = self._extract_metadata(file_path)
            file_id = metadata.file_id
            
            # Check if file has changed since last index
            if self._file_unchanged(file_id, metadata.modified_at):
                logger.debug(f"File unchanged: {name}")
                return ProcessingStatus.SKIPPED
            
            # Extract content
            content = self._extract_content(file_path, metadata.mime_type)
            
            # Index metadata
            self.metadata_collection.upsert(
                documents=[str(metadata)],
                metadatas=[asdict(metadata)],
                ids=[f"meta-{file_id}"]
            )
            
            # Index content if available
            if content:
                self.content_collection.upsert(
                    documents=[content],
                    metadatas=[asdict(metadata)],
                    ids=[f"content-{file_id}"]
                )
            
            logger.debug(f"Indexed file: {name}")
            return ProcessingStatus.SUCCESS
            
        except Exception as e:
            logger.error(f"Error indexing file {file_path}: {e}")
            return ProcessingStatus.FAILURE
    
    def _delete_file(self, file_path: str) -> bool:
        """Delete a file from the index."""
        try:
            file_id = self._get_file_hash(file_path)
            
            # Check if the file exists in metadata collection before deleting
            try:
                existing = self.metadata_collection.get(ids=[f"meta-{file_id}"])
                if existing["ids"]:
                    # Delete from both collections only if it exists
                    self.metadata_collection.delete(ids=[f"meta-{file_id}"])
                    self.content_collection.delete(ids=[f"content-{file_id}"])
                    logger.debug(f"Deleted file from index: {file_path}")
                else:
                    logger.debug(f"File not in index, skipping deletion: {file_path}")
            except Exception:
                # If get() fails, try to delete anyway (might exist in content but not metadata)
                self.metadata_collection.delete(ids=[f"meta-{file_id}"])
                self.content_collection.delete(ids=[f"content-{file_id}"])
                logger.debug(f"Attempted deletion of file from index: {file_path}")
            
            return True
            
        except Exception as e:
            logger.error(f"Error deleting file {file_path}: {e}")
            return False
    
    def _move_file(self, task: FileTask) -> bool:
        """Handle file move operation (atomic delete + index)."""
        if not task.metadata:
            logger.error("Move task missing metadata")
            return False
        
        old_path = task.metadata.get('old_path')
        new_path = task.metadata.get('new_path')
        
        if not old_path or not new_path:
            logger.error("Move task missing old_path or new_path")
            return False
        
        # Atomic operation: delete old, then index new
        delete_success = self._delete_file(old_path)
        index_success = self._index_file(new_path)
        
        if delete_success and index_success:
            logger.debug(f"Moved file: {old_path} -> {new_path}")
            return True
        else:
            logger.error(f"Failed to move file: {old_path} -> {new_path}")
            return False
    
    def _get_file_hash(self, file_path: str) -> str:
        """Generate hash for file path."""
        return hashlib.sha256(file_path.encode()).hexdigest()
    
    def _extract_metadata(self, file_path: str) -> FileMetadata:
        """Extract metadata from file."""
        stat = os.stat(file_path)
        path = Path(file_path)
        
        mime_type, _ = mimetypes.guess_type(file_path)
        if mime_type is None:
            mime_type = "application/octet-stream"
        
        return FileMetadata(
            file_id=self._get_file_hash(file_path),
            name=path.name,
            extension=path.suffix.lower(),
            path=str(path),
            parent_dir=str(path.parent),
            size=stat.st_size,
            created_at=datetime.fromtimestamp(stat.st_ctime).isoformat(),
            modified_at=datetime.fromtimestamp(stat.st_mtime).isoformat(),
            accessed_at=datetime.fromtimestamp(stat.st_atime).isoformat(),
            mime_type=mime_type,
        )
    
    def _check_file_size(self, file_path: str) -> bool:
        """Check if file size is within limits."""
        file_size = os.path.getsize(file_path)
        file_extension = Path(file_path).suffix.lower()
        
        # Different size limits for different file types
        size_limits = {
            'pdf': 50_000_000,      # 50MB
            'office': 20_000_000,   # 20MB
            'text': 5_000_000,      # 5MB
            'default': 10_000_000   # 10MB
        }
        
        if file_extension == ".pdf":
            size_limit = size_limits['pdf']
        elif file_extension in [".docx", ".xlsx", ".pptx"]:
            size_limit = size_limits['office']
        elif file_extension in [".txt", ".md", ".py", ".js", ".html", ".css"]:
            size_limit = size_limits['text']
        else:
            size_limit = size_limits['default']
        
        if file_size > size_limit:
            logger.debug(f"Skipping large file: {os.path.basename(file_path)} ({file_size:,} bytes)")
            return False
        
        return True
    
    def _file_unchanged(self, file_id: str, modified_at: str) -> bool:
        """Check if file has changed since last index."""
        try:
            result = self.metadata_collection.get(where={"file_id": file_id})
            if result["metadatas"]:
                existing_metadata = result["metadatas"][0]
                return existing_metadata["modified_at"] == modified_at
        except Exception:
            pass
        return False
    
    def _extract_content(self, file_path: str, mime_type: str) -> Optional[str]:
        """Extract content from various file types."""
        file_extension = Path(file_path).suffix.lower()
        
        try:
            # PDF files
            if mime_type == "application/pdf" or file_extension == ".pdf":
                return self._extract_pdf_content(file_path)
            
            # Microsoft Office files
            elif file_extension == ".docx":
                return self._extract_docx_content(file_path)
            elif file_extension in [".xlsx", ".xlsm"]:
                return self._extract_excel_content(file_path)
            elif file_extension == ".pptx":
                return self._extract_pptx_content(file_path)
            
            # HTML and Markdown
            elif file_extension in [".html", ".htm", ".xhtml"]:
                return self._extract_html_content(file_path)
            elif file_extension in [".md", ".markdown"]:
                return self._extract_markdown_content(file_path)
            
            # Text files
            elif (mime_type and mime_type.startswith("text/")) or file_extension in [
                ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".cpp", ".c", ".h",
                ".cs", ".php", ".rb", ".go", ".rs", ".swift", ".kt", ".scala",
                ".css", ".scss", ".sass", ".less", ".json", ".xml", ".yaml", ".yml",
                ".ini", ".cfg", ".conf", ".log", ".sql", ".sh", ".bash", ".zsh",
                ".txt", ".csv", ".tsv"
            ]:
                return self._extract_text_content(file_path)
            
            else:
                logger.debug(f"Unsupported file type: {mime_type} for {file_path}")
                return None
                
        except Exception as e:
            logger.error(f"Error extracting content from {file_path}: {e}")
            return None
    
    def _detect_encoding(self, file_path: str) -> str:
        """Detect file encoding using chardet."""
        with open(file_path, 'rb') as f:
            raw_data = f.read(32 * 1024)  # Read first 32KB
            result = chardet.detect(raw_data)
            return result['encoding'] or 'utf-8'
    
    def _extract_text_content(self, file_path: str) -> Optional[str]:
        """Extract content from text files."""
        encoding = self._detect_encoding(file_path)
        try:
            with open(file_path, 'r', encoding=encoding, errors='strict') as f:
                return f.read()
        except (UnicodeDecodeError, TypeError, LookupError):
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"Error reading text file {file_path}: {e}")
                return None
    
    def _extract_pdf_content(self, file_path: str) -> Optional[str]:
        """Extract text from PDF files."""
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PDF content: {e}")
            return None
    
    def _extract_docx_content(self, file_path: str) -> Optional[str]:
        """Extract text from DOCX files."""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting DOCX content: {e}")
            return None
    
    def _extract_excel_content(self, file_path: str) -> Optional[str]:
        """Extract text from Excel files."""
        try:
            workbook = load_workbook(file_path, data_only=True)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting Excel content: {e}")
            return None
    
    def _extract_pptx_content(self, file_path: str) -> Optional[str]:
        """Extract text from PowerPoint files."""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PPTX content: {e}")
            return None
    
    def _extract_html_content(self, file_path: str) -> Optional[str]:
        """Extract text from HTML files."""
        html_content = self._extract_text_content(file_path)
        if not html_content:
            return None
        
        try:
            soup = BeautifulSoup(html_content, 'html.parser')
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            return text
        except Exception as e:
            logger.error(f"Error parsing HTML content: {e}")
            return None
    
    def _extract_markdown_content(self, file_path: str) -> Optional[str]:
        """Extract text from Markdown files."""
        md_content = self._extract_text_content(file_path)
        if not md_content:
            return None
        
        try:
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            return soup.get_text()
        except Exception as e:
            logger.error(f"Error parsing Markdown content: {e}")
            return None