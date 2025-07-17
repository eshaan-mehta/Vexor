import os
import hashlib
import mimetypes
import chromadb
import threading
import time
import json
import atexit
import signal
import sys
from concurrent.futures import ThreadPoolExecutor, as_completed

from datetime import datetime
from dataclasses import asdict
from typing import Dict, Any, List, Tuple, Callable, Optional
from chromadb.config import Settings
from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction
from sentence_transformers import SentenceTransformer
from pathlib import Path

# New imports for file type support
import chardet
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
import markdown

from models.filemetadata import FileMetadata

class Indexer:
    embedding_function = SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")  

    def __init__(self, 
            db_path: str = "./chroma", 
            metadata_collection_name: str = "file_metadata",
            content_collection_name: str = "file_content",
            batch_size: int = 50
        ):
        os.makedirs(db_path, exist_ok=True)
        
        self.db_path = db_path
        self.batch_size = batch_size
        self.pending_metadata = [] 
        self.pending_content = []
        self.batch_backup_file = os.path.join(db_path, "batch_backup.json")
        
        # Register cleanup handlers for graceful shutdown
        atexit.register(self._emergency_flush)
        signal.signal(signal.SIGINT, self._signal_handler)
        signal.signal(signal.SIGTERM, self._signal_handler)
        
        # Check for and recover any pending batches from previous crashes
        self._recover_from_crash()

        self.client = chromadb.PersistentClient(
            path=db_path,
            settings=Settings(anonymized_telemetry=False)
        )

        self.content_collection = self.client.get_or_create_collection(
            name=content_collection_name,
            embedding_function=self.embedding_function,
            metadata = {
                "hnsw:space": "cosine",  # cosine better for text embeddings
                "hnsw:construction_ef": 100,  # size of candidates during indexing (default = 100)
                "hnsw:search_ef": 100,  # size of candidates during searching (default = 100)
                "hnsw:M": 16  # max neighbours in node graph (default = 16)
            }
        )

        self.metadata_collection = self.client.get_or_create_collection(
            name=metadata_collection_name,
            embedding_function=self.embedding_function,
            metadata = {
                "hnsw:space": "cosine",  # cosine better for text embeddings
                "hnsw:construction_ef": 100,  # size of candidates during indexing (default = 100)
                "hnsw:search_ef": 100,  # size of candidates during searching (default = 100)
                "hnsw:M": 16  # max neighbours in node graph (default = 16)
            }
        )
        
    def get_file_hash(self, file_path: str) -> str:
        # TODO: Update to use file contents
        return  hashlib.sha256(file_path.encode()).hexdigest()

    def __extract_metadata(self, file_path: str) -> FileMetadata:
        stat = os.stat(file_path)
        path = Path(file_path)
        
        mime_type, _ = mimetypes.guess_type(file_path)
        if mime_type is None:
            mime_type = "application/octet-stream"  # Default binary type

        return FileMetadata(
            file_id=self.get_file_hash(file_path),
            name=path.name,
            extension=path.suffix.lower(),
            path=str(path),
            parent_dir=str(path.parent),
            size=stat.st_size,
            created_at=datetime.fromtimestamp(stat.st_ctime).isoformat(),
            modified_at=datetime.fromtimestamp(stat.st_mtime).isoformat(),
            accessed_at=datetime.fromtimestamp(stat.st_atime).isoformat(),
            mime_type=mime_type,
        )

    def __detect_encoding(self, file_path: str) -> str:
        """Detect file encoding using chardet"""
        with open(file_path, 'rb') as f:
            raw_data = f.read(32 * 1024) # Read first 32KB for speed
            result = chardet.detect(raw_data)
            return result['encoding'] or 'utf-8'

    def __read_file(self, file_path: str) -> str | None:
        """Reads a file with robust encoding detection and fallback."""
        encoding = self.__detect_encoding(file_path)
        try:
            with open(file_path, 'r', encoding=encoding, errors='strict') as f:
                return f.read()
        except (UnicodeDecodeError, TypeError, LookupError):
            try:
                # Fallback to utf-8 with error handling
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception as e:
                print(f"Error reading file {file_path} after fallback: {e}")
                return None

    def __extract_text_content(self, file_path: str) -> str:
        """Extract content from text files with proper encoding detection"""
        return self.__read_file(file_path)

    def __extract_pdf_content(self, file_path: str) -> str:
        """Extract text content from PDF files"""
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting PDF content: {e}")
            return None

    def __extract_docx_content(self, file_path: str) -> str:
        """Extract text content from DOCX files"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting DOCX content: {e}")
            return None

    def __extract_excel_content(self, file_path: str) -> str:
        """Extract text content from Excel files"""
        try:
            workbook = load_workbook(file_path, data_only=True)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting Excel content: {e}")
            return None

    def __extract_pptx_content(self, file_path: str) -> str:
        """Extract text content from PowerPoint files"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text.strip()
        except Exception as e:
            print(f"Error extracting PPTX content: {e}")
            return None

    def __extract_html_content(self, file_path: str) -> str:
        """Extract text content from HTML files"""
        html_content = self.__read_file(file_path)
        if not html_content:
            return None
        
        try:
            soup = BeautifulSoup(html_content, 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text()
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            return text
        except Exception as e:
            print(f"Error parsing HTML content: {e}")
            return None

    def __extract_markdown_content(self, file_path: str) -> str:
        """Extract text content from Markdown files"""
        md_content = self.__read_file(file_path)
        if not md_content:
            return None
        
        try:
            # Convert markdown to HTML then extract text
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            return soup.get_text()
        except Exception as e:
            print(f"Error parsing Markdown content: {e}")
            return None

    def __extract_content(self, file_path: str, mime_type: str) -> str:
        """Extract content from various file types"""
        file_extension = Path(file_path).suffix.lower()
        
        # PDF files
        if mime_type == "application/pdf" or file_extension == ".pdf":
            return self.__extract_pdf_content(file_path)
        
        # Microsoft Office files
        elif mime_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document"] or file_extension == ".docx":
            return self.__extract_docx_content(file_path)
        
        elif mime_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"] or file_extension in [".xlsx", ".xlsm"]:
            return self.__extract_excel_content(file_path)
        
        elif mime_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation"] or file_extension == ".pptx":
            return self.__extract_pptx_content(file_path)
        
        # HTML and XML files
        elif mime_type in ["text/html", "application/xhtml+xml"] or file_extension in [".html", ".htm", ".xhtml"]:
            return self.__extract_html_content(file_path)
        
        # Markdown files
        elif mime_type == "text/markdown" or file_extension in [".md", ".markdown"]:
            return self.__extract_markdown_content(file_path)
        
        # Programming language files and other text files
        elif (mime_type and mime_type.startswith("text/")) or file_extension in [
            ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".cpp", ".c", ".h", 
            ".cs", ".php", ".rb", ".go", ".rs", ".swift", ".kt", ".scala",
            ".css", ".scss", ".sass", ".less", ".json", ".xml", ".yaml", ".yml",
            ".ini", ".cfg", ".conf", ".log", ".sql", ".sh", ".bash", ".zsh",
            ".txt", ".csv", ".tsv"
        ]:
            return self.__extract_text_content(file_path)
        
        else:
            print(f"Unsupported file type: {mime_type} for {file_path}")
            return None

    def index_file(self, file_path: str, use_batch: bool = False) -> bool:
        if os.path.isdir(file_path) or not os.path.exists(file_path):
            return False
        
        name = os.path.basename(file_path)
        
         # skip hidden files and temporary office files
        if name.startswith((".", "__", "~$")):
            print(f"\nSkipping hidden/temporary file: {name}")
            return False
            
        # Skip large files until TODO: chunking (different limits for different file types)
        file_size = os.path.getsize(file_path)
        file_extension = Path(file_path).suffix.lower()
        
        # Different size limits for different file types
        size_limits = {
            # Text files - smaller limit
            'text': 5_000_000,  # 5MB
            # Office documents - larger limit as they often contain more content
            'office': 20_000_000,  # 20MB
            # PDFs - larger limit
            'pdf': 50_000_000,  # 50MB
            # Default
            'default': 10_000_000  # 10MB
        }
        
        if file_extension == ".pdf":
            size_limit = size_limits['pdf']
        elif file_extension in [".docx", ".xlsx", ".pptx"]:
            size_limit = size_limits['office']
        elif file_extension in [".txt", ".md", ".py", ".js", ".html", ".css"]:
            size_limit = size_limits['text']
        else:
            size_limit = size_limits['default']
            
        if file_size > size_limit:
            print(f"\nSkipping large file: {name} ({file_size:,} bytes > {size_limit:,} limit)")
            return False
        
        try:
            metadata = self.__extract_metadata(file_path)
            file_id = metadata.file_id

            # Check if file has changed since last index
            if result := self.metadata_collection.get(where={"file_id": file_id}):
                if result["metadatas"]:
                    existing_metadata = result["metadatas"][0]
                    if existing_metadata["modified_at"] == metadata.modified_at:
                        print(f"No change in: {name}")
                        return False

            content = self.__extract_content(file_path, metadata.mime_type)
            
            if use_batch:
                # Add to batch for later processing
                self._add_to_batch(metadata, content)
                
                # Check if we should flush the batch
                if len(self.pending_metadata) >= self.batch_size:
                    self._flush_batch()
            else:
                # Process immediately (for single file operations)
                self.metadata_collection.upsert(
                    documents=[str(metadata)],
                    metadatas=[asdict(metadata)],
                    ids=[f"meta-{file_id}"]
                )

                if content:
                    self.content_collection.upsert(
                        documents=[content],
                        metadatas=[asdict(metadata)],
                        ids=[f"content-{file_id}"]
                    )

            # split file into overlapping chunks, will have a seperate entry for each
            # generate hash for each chunk
            # check if hash has changed since last index
            # multi thread this part to index chunks in parallel
            # update db with new chunks
            print(f"Done Indexing file: {name}")
            return True
        except Exception as e:
            print(f"Could not index file: {name}")
            print(f"Error: {e}")
            return False

    def delete_file(self, file_path: str) -> bool:
        """Deletes a file's entries from the database based on its path."""
        try:
            file_id = self.get_file_hash(file_path)
            # Deleting from both collections. ChromaDB handles non-existent IDs gracefully.
            self.metadata_collection.delete(ids=[f"meta-{file_id}"])
            self.content_collection.delete(ids=[f"content-{file_id}"])
            print(f"Deleted indexed data for: {file_path}")
            return True
        except Exception as e:
            print(f"Error deleting file {file_path}: {e}")
            return False

    def cleanup_deleted_files(self) -> int:
        """Remove entries for files that no longer exist on disk."""
        deleted_count = 0
        
        try:
            # Get all metadata entries
            all_metadata = self.metadata_collection.get()
            
            if not all_metadata['metadatas']:
                print("No files in database to check")
                return 0
            
            print(f"Checking {len(all_metadata['metadatas'])} files for deletion...")
            
            for metadata in all_metadata['metadatas']:
                file_path = metadata.get('path')
                if file_path and not os.path.exists(file_path):
                    print(f"File no longer exists: {file_path}")
                    if self.delete_file(file_path):
                        deleted_count += 1
            
            if deleted_count > 0:
                print(f"Cleaned up {deleted_count} deleted files from database")
            else:
                print("No deleted files found in database")
                
        except Exception as e:
            print(f"Error during cleanup: {e}")
        
        return deleted_count

    def _add_to_batch(self, metadata: FileMetadata, content: str = None):
        """Add items to batch for later processing"""
        file_id = metadata.file_id
        
        # Add metadata to batch
        self.pending_metadata.append({
            'document': str(metadata),
            'metadata': asdict(metadata),
            'id': f"meta-{file_id}"
        })
        
        # Add content to batch if available
        if content:
            self.pending_content.append({
                'document': content,
                'metadata': asdict(metadata),
                'id': f"content-{file_id}"
            })
        
        # Backup batch to disk for crash recovery
        self._backup_batch()

    def _flush_batch(self):
        """Process all pending batch items"""
        if self.pending_metadata:
            print(f"Flushing {len(self.pending_metadata)} metadata items...")
            
            # Prepare batch data
            documents = [item['document'] for item in self.pending_metadata]
            metadatas = [item['metadata'] for item in self.pending_metadata]
            ids = [item['id'] for item in self.pending_metadata]
            
            # Batch upsert metadata
            self.metadata_collection.upsert(
                documents=documents,
                metadatas=metadatas,
                ids=ids
            )
            self.pending_metadata.clear()
        
        if self.pending_content:
            print(f"Flushing {len(self.pending_content)} content items...")
            
            # Prepare batch data
            documents = [item['document'] for item in self.pending_content]
            metadatas = [item['metadata'] for item in self.pending_content]
            ids = [item['id'] for item in self.pending_content]
            
            # Batch upsert content
            self.content_collection.upsert(
                documents=documents,
                metadatas=metadatas,
                ids=ids
            )
            self.pending_content.clear()

    def index_directory(self, root_dir: str, cleanup_deleted: bool = True, use_batch: bool = False) -> int:
        if not os.path.exists(root_dir):
            raise FileNotFoundError(f"Directory not found: {root_dir}")

        # First, clean up deleted files if requested
        if cleanup_deleted:
            print("Checking for deleted files...")
            deleted_count = self.cleanup_deleted_files()
            if deleted_count > 0:
                print(f"Removed {deleted_count} deleted files from database")

        # Recursively loop through all subdirectories using batch processing
        count = 0
        print(f"Starting directory indexing for: {root_dir}")
        
        try:
            for path, dirs, files in os.walk(root_dir):
                for file in files:
                    file_path = os.path.join(path, file)
                    try:
                        success = self.index_file(file_path, use_batch=use_batch)
                        if success:
                            count += 1
                        
                    except Exception as e:
                        print(f"Exception indexing {file_path}: {e}")
        finally:
            # Always flush remaining batch items at the end
            if use_batch:
                print("Flushing remaining batch items...")
                self._flush_batch()
        
        print(f"Directory indexing complete. Indexed {count} files from {root_dir}")
        return count

    def _backup_batch(self):
        """Backup current batch to disk for crash recovery"""
        if not self.pending_metadata and not self.pending_content:
            # No pending items, remove backup file if it exists
            if os.path.exists(self.batch_backup_file):
                os.remove(self.batch_backup_file)
            return
            
        try:
            backup_data = {
                'timestamp': datetime.now().isoformat(),
                'pending_metadata': self.pending_metadata,
                'pending_content': self.pending_content
            }
            
            with open(self.batch_backup_file, 'w') as f:
                json.dump(backup_data, f, indent=2)
                
        except Exception as e:
            print(f"Warning: Could not backup batch data: {e}")

    def _recover_from_crash(self):
        """Recover and process any pending batches from previous crashes"""
        if not os.path.exists(self.batch_backup_file):
            return
            
        try:
            print("Found batch backup file from previous session...")
            
            with open(self.batch_backup_file, 'r') as f:
                backup_data = json.load(f)
            
            # Restore pending batches
            self.pending_metadata = backup_data.get('pending_metadata', [])
            self.pending_content = backup_data.get('pending_content', [])
            
            if self.pending_metadata or self.pending_content:
                print(f"Recovering {len(self.pending_metadata)} metadata and {len(self.pending_content)} content items...")
                
                # Process the recovered batch
                self._flush_batch()
                print("Successfully recovered and processed pending batch items")
            
            # Clean up backup file
            os.remove(self.batch_backup_file)
            
        except Exception as e:
            print(f"Error recovering from crash: {e}")
            # If recovery fails, clear the backup file to prevent repeated errors
            try:
                os.remove(self.batch_backup_file)
            except:
                pass

    def _emergency_flush(self):
        """Emergency flush called during shutdown"""
        try:
            if self.pending_metadata or self.pending_content:
                print("Emergency flush: Processing pending batch items...")
                self._flush_batch()
                
                # Clean up backup file after successful flush
                if os.path.exists(self.batch_backup_file):
                    os.remove(self.batch_backup_file)
                    
        except Exception as e:
            print(f"Error during emergency flush: {e}")

    def _signal_handler(self, signum, frame):
        """Handle shutdown signals gracefully"""
        print(f"\nReceived signal {signum}, performing emergency flush...")
        self._emergency_flush()
        sys.exit(0)